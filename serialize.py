import sys
import json
from pptx import Presentation
from pptx.shapes.picture import Picture
from pptx.shapes.autoshape import Shape
from pptx.shapes.group import GroupShape
from pptx.util import Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

def escape(text):
    return text.replace('\x0b', '\n')

def shape_repr(shape, slide_size):
    if isinstance(shape, Shape) and shape.text:
        return [escape(shape.text)]
    elif isinstance(shape, GroupShape):
        return shapes_repr(shape.shapes, slide_size)
    else:
        return []

def shapes_repr(shapes, slide_size):
    return sum([shape_repr(shape, slide_size) for shape in shapes], [])

def slide_repr(slide, slide_size):
    return "\n".join(shapes_repr(slide.shapes, slide_size))

def main(filename):
    prs = Presentation(filename)
    slide_size = (prs.slide_width, prs.slide_height)
    slide_reprs = [slide_repr(slide, slide_size) for slide in prs.slides]
    slides_repr = "\n\n---\n\n".join(slide_reprs)
    print(slides_repr)


if __name__ == '__main__':
    main(sys.argv[1])

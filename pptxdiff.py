import sys
import os 
from pptx import Presentation
from serialize import shape_repr
from lcs import LCS
def head(ls):
    return ls[0] if ls else None

def slide_repr(slide, slide_size):
    return list(filter(None, [shape_repr(shape, slide_size) for shape in slide.shapes]))

def labels(length, idx, head):
    if length > 0:
        return [head + str(idx + i) for i in range(length)]
    else:
        return [head + str(idx)]

def matchings(lines1, lines2):
    lcs_idxs = []
    for i in range(len(lines1)):
        for j in range(len(lines2)):
            lcs_idxs.append(((i, j), LCS(lines1[i], lines2[j])))
    list.sort(lcs_idxs, key=lambda li: li[1].size)
    ranges1 = [(0, len(lines2) - 1) for _ in range(len(lines1))]
    ranges2 = [(0, len(lines1) - 1) for _ in range(len(lines2))]
    matches1 = [None for _ in range(len(lines1))]
    matches2 = [None for _ in range(len(lines2))]
    while len(lcs_idxs) > 0:
        (i, j), lcs = lcs_idxs.pop()
        if matches1[i] is not None or matches2[j] is not None:
            continue
        if ranges1[i][0] > j or ranges1[i][1] < j or ranges2[j][0] > i or ranges2[j][1] < i:
            continue
        matches1[i] = (j, lcs)
        matches2[j] = (i, lcs)
        for i_ in range(len(lines1)):
            if i_ < i:
                ranges1[i_] = (ranges1[i_][0], min(ranges1[i_][1], j - 1))
            if i_ > i:
                ranges1[i_] = (max(ranges1[i_][0], j + 1), ranges1[i_][1])
        for j_ in range(len(lines2)):
            if j_ < j:
                ranges2[j_] = (ranges2[j_][0], min(ranges2[j_][1], i - 1))
            if j_ > j:
                ranges2[j_] = (max(ranges2[j_][0], i + 1), ranges2[j_][1])
    idx1 = 0
    idx2 = 0
    result = []
    while True:
        while idx1 < len(matches1) and matches1[idx1] is None:
            result.append((idx1, None, None))
            idx1 += 1
        while idx2 < len(matches2) and matches2[idx2] is None:
            result.append((None, idx2, None))
            idx2 += 1
        if idx1 == len(matches1) or idx2 == len(matches2):
            break
        result.append((idx1, idx2, matches1[idx1][1]))
        idx1 += 1
        idx2 += 1
    return result

def show_diff(lines1, lines2, idx1, idx2):
    m = matchings(lines1, lines2)
    for i, j, lcs in m:
        label1 = "Slide" + str(idx1 + 1)
        label2 = "Slide" + str(idx2 + 1)
        if i is None:
            els1 = ""
            els2 = "\n".join(lines2[j])
            idx2 += 1
        elif j is None:
            els1 = "\n".join(lines1[i])
            els2 = ""
            idx1 += 1
        else:
            diff1, diff2 = lcs.diff_collection()
            els1 = "\n".join([lines1[i][n] for n in diff1])
            els2 = "\n".join([lines2[j][m] for m in diff2])
            idx1 += 1
            idx2 += 1
        if els1 == "":
            print(label1 + " a " + label2)
        elif els2 == "":
            print(label1 + " d " + label2)
        else:
            print(label1 + " c " + label2)
        print("<<<")
        print(els1)
        print(">>>")
        print(els2)
        print()
    

def main(file1, file2):
    prs1 = Presentation(file1)
    slide_size = (prs1.slide_width, prs1.slide_height)
    slides1 = [slide_repr(slide, slide_size) for slide in prs1.slides]
    prs2 = Presentation(file2)
    slides2 = [slide_repr(slide, slide_size) for slide in prs2.slides]
    lcs = LCS(slides1, slides2)
    for ((l1, i1), (l2, i2)) in lcs.diff():
        lines1 = [slides1[i + i1] for i in range(l1)]
        lines2 = [slides2[i + i2] for i in range(l2)]
        show_diff(lines1, lines2, i1, i2)

if __name__ == '__main__':
    main(sys.argv[1], sys.argv[2])
